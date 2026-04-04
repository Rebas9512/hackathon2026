#!/usr/bin/env python3
"""Generate opening slides as PPTX for Spillover Alpha presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
CARD = RGBColor(0x11, 0x11, 0x14)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x71, 0x71, 0x7A)
SECONDARY = RGBColor(0xA1, 0xA1, 0xAA)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
BORDER = RGBColor(0x27, 0x27, 0x2A)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_name="Inter",
             font_size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_name="Inter",
                  font_size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                  line_spacing=1.2):
    """Add text box with multiple styled lines."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, lcolor, lbold, lsize) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(lsize if lsize else font_size)
        p.font.color.rgb = lcolor if lcolor else color
        p.font.bold = lbold if lbold is not None else bold
        p.alignment = alignment
        p.line_spacing = Pt((lsize if lsize else font_size) * line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide1, BLACK)

# Tag
add_text(slide1, 1.0, 0.8, 14, 0.4,
         "2026 FINHACK CHALLENGE  ·  CASE 4  ·  TEAM 54",
         "Courier New", 14, PURPLE, bold=True)

# Title
add_multiline(slide1, 1.0, 1.5, 10, 3, [
    ("Spillover", WHITE, True, 96),
    ("Alpha", WHITE, True, 96),
], font_name="Georgia", line_spacing=0.95)

# Subtitle
add_multiline(slide1, 1.0, 4.2, 8, 1.5, [
    ("Can cross-company sentiment contagion", SECONDARY, False, 30),
    ("predict post-earnings stock returns?", SECONDARY, False, 30),
], font_name="Inter", line_spacing=1.3)

# Accent line
add_rect(slide1, 1.0, 5.5, 0.6, 0.03, PURPLE)

# Meta
add_multiline(slide1, 1.0, 5.7, 10, 1, [
    ("Multi-Source Sentiment Spillover Network × Post-Earnings Prediction", MUTED, False, 18),
    ("Magnificent 7: AAPL · MSFT · GOOGL · AMZN · NVDA · META · TSLA", SECONDARY, False, 18),
    ("91 Earnings Events  ·  2022 Q1 – 2025 Q1  ·  5 Progressive Models", MUTED, False, 18),
], font_name="Inter", line_spacing=1.5)

# Badge
badge = add_rect(slide1, 1.0, 6.8, 2.3, 0.45, BLACK, PURPLE)
add_text(slide1, 1.0, 6.82, 2.3, 0.45,
         "UTD JSOM Finance Lab",
         "Courier New", 12, PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Right accent gradient (just a subtle purple rect)
add_rect(slide1, 12.5, 0, 3.5, 9, RGBColor(0x12, 0x0C, 0x1A))


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Hook
# ═══════════════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2, BLACK)

# Left: The paradox
add_multiline(slide2, 1.0, 1.5, 7, 2, [
    ("Mag7 beats earnings", WHITE, True, 52),
    ("90% of the time.", WHITE, True, 52),
], font_name="Georgia", line_spacing=1.05)

add_multiline(slide2, 1.0, 3.5, 7, 2, [
    ("The stock drops", RED, True, 52),
    ("nearly half the time.", RED, True, 52),
], font_name="Georgia", line_spacing=1.05)

# Accent line
add_rect(slide2, 1.0, 5.5, 0.6, 0.03, PURPLE)

# Hook question
add_multiline(slide2, 1.0, 5.7, 7, 1.2, [
    ("What if the answer is hiding in the network", SECONDARY, False, 24),
    ("between companies — not inside any single one?", SECONDARY, False, 24),
], font_name="Inter", line_spacing=1.4)

# Right: Stats card
card_left = 9.2
card_top = 1.2
card_w = 5.8
card_h = 6.5
add_rect(slide2, card_left, card_top, card_w, card_h, CARD, BORDER)

# 90% row
add_text(slide2, card_left + 0.5, card_top + 0.5, 2.5, 0.8,
         "90%", "Georgia", 52, GREEN, bold=True)
add_multiline(slide2, card_left + 2.8, card_top + 0.55, 2.5, 0.8, [
    ("EPS beat rate", MUTED, False, 18),
    ("(Mag7, 2022–2025)", MUTED, False, 18),
], line_spacing=1.3)

# Divider
add_rect(slide2, card_left + 0.5, card_top + 1.7, card_w - 1, 0.01, BORDER)

# 49% row
add_text(slide2, card_left + 0.5, card_top + 2.0, 2.5, 0.8,
         "49%", "Georgia", 52, RED, bold=True)
add_multiline(slide2, card_left + 2.8, card_top + 2.05, 2.5, 0.8, [
    ("Stock drops anyway", MUTED, False, 18),
    ("after beating EPS", MUTED, False, 18),
], line_spacing=1.3)

# Divider
add_rect(slide2, card_left + 0.5, card_top + 3.2, card_w - 1, 0.01, BORDER)

# +205% row
add_text(slide2, card_left + 0.5, card_top + 3.5, 2.5, 0.8,
         "+205%", "Georgia", 52, YELLOW, bold=True)
add_multiline(slide2, card_left + 2.8, card_top + 3.55, 2.5, 0.8, [
    ("Our best strategy return", MUTED, False, 18),
    ("(spillover network + XGBoost)", MUTED, False, 18),
], line_spacing=1.3)

# Divider
add_rect(slide2, card_left + 0.5, card_top + 4.7, card_w - 1, 0.01, BORDER)

# CTA
add_text(slide2, card_left + 0.5, card_top + 5.1, 4, 0.5,
         "Let us show you how →",
         "Inter", 20, PURPLE, bold=True)


# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════

out_path = os.path.join(os.path.dirname(__file__), "Spillover_Alpha_Opening.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
